"""
组会演示文稿生成器
论文大纲 + 六工况前沿分析成果展示
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import os

# ── 路径 ─────────────────────────────────────────────────────────────────────
PLOT_DIR = '/Users/a1234/Carnot_Battery_for_DC/pure_deap_nsga/plots/global_pareto'
OUT_PATH = '/Users/a1234/Carnot_Battery_for_DC/pure_deap_nsga/reports/group_meeting_presentation.pptx'

# ── 颜色主题 ─────────────────────────────────────────────────────────────────
DARK_BLUE  = RGBColor(0x1A, 0x3A, 0x5C)   # 深蓝 — 标题
MID_BLUE   = RGBColor(0x21, 0x66, 0xAC)   # 中蓝 — 重点
ACCENT_RED = RGBColor(0xCC, 0x33, 0x00)   # 红   — 警示/创新
ACCENT_GRN = RGBColor(0x1A, 0x96, 0x41)   # 绿   — 完成
LIGHT_GRAY = RGBColor(0xF2, 0xF4, 0xF7)   # 浅灰 — 背景框
MID_GRAY   = RGBColor(0x88, 0x88, 0x88)   # 灰   — 次要文字
BLACK      = RGBColor(0x1A, 0x1A, 0x1A)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
YELLOW_HL  = RGBColor(0xFF, 0xF0, 0xC0)   # 高亮黄

# ── 幻灯片尺寸（16:9 宽屏）─────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height
BLANK = prs.slide_layouts[6]   # 全空白

# ── 辅助函数 ─────────────────────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid(); shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        if line_w: shape.line.width = Pt(line_w)
    else:
        shape.line.fill.background()
    return shape

def add_txt(slide, text, l, t, w, h,
            size=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT,
            wrap=True, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def add_img(slide, path, l, t, w, h=None):
    if not os.path.exists(path):
        print(f'  [WARN] image not found: {path}')
        return None
    if h:
        return slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        return slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w))

def header_bar(slide, title, subtitle=None):
    """顶部蓝色横幅"""
    add_rect(slide, 0, 0, 13.33, 1.05, fill=DARK_BLUE)
    add_txt(slide, title, 0.35, 0.08, 10, 0.6,
            size=24, bold=True, color=WHITE)
    if subtitle:
        add_txt(slide, subtitle, 0.35, 0.62, 11, 0.38,
                size=13, color=RGBColor(0xB8, 0xD4, 0xF0))

def slide_num(slide, n, total=14):
    add_txt(slide, f'{n} / {total}', 12.1, 7.1, 1.1, 0.35,
            size=10, color=MID_GRAY, align=PP_ALIGN.RIGHT)

def bullet(slide, items, l, t, w, h, size=15, indent=0.3, spacing=0.42,
           color=BLACK, marker='▸', bold_first=False):
    """简单项目符号列表"""
    for i, item in enumerate(items):
        is_sub = item.startswith('  ')
        text = (marker + '  ' if not is_sub else '    –  ') + item.strip()
        fc = MID_GRAY if is_sub else color
        fs = size - 1 if is_sub else size
        add_txt(slide, text, l + (0.3 if is_sub else 0), t + i * spacing,
                w, 0.4, size=fs, color=fc, bold=(bold_first and i == 0))

def badge(slide, text, l, t, color=MID_BLUE, text_color=WHITE, w=1.6, h=0.32):
    add_rect(slide, l, t, w, h, fill=color)
    add_txt(slide, text, l + 0.05, t + 0.02, w - 0.1, h,
            size=11, bold=True, color=text_color, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — 封面
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)

# 深蓝渐变背景条
add_rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xF5, 0xF8, 0xFC))
add_rect(sl, 0, 0, 13.33, 2.8, fill=DARK_BLUE)

add_txt(sl, 'Multi-Objective Optimization of Carnot Battery Systems',
        0.5, 0.35, 12.3, 0.8,
        size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_txt(sl, 'for Data Center Waste Heat Recovery',
        0.5, 1.1, 12.3, 0.6,
        size=22, bold=False, color=RGBColor(0xB8, 0xD4, 0xF0), align=PP_ALIGN.CENTER)

add_rect(sl, 1.2, 2.5, 10.9, 0.05, fill=RGBColor(0x4A, 0x90, 0xD9))

add_txt(sl, '论文进度报告 · 六工况全局帕累托前沿分析',
        0.5, 2.75, 12.3, 0.55,
        size=17, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)

# 信息框
for i, (label, val) in enumerate([
    ('日期', '2026 年 4 月'),
    ('阶段', 'Part I 静态性能映射'),
    ('工况', 'DC-A → DC-F  共六个'),
    ('目标', '三目标 NSGA-II 多构型优化'),
]):
    x = 0.6 + i * 3.1
    add_rect(sl, x, 3.55, 2.9, 0.95, fill=WHITE, line=MID_BLUE, line_w=0.8)
    add_txt(sl, label, x + 0.12, 3.60, 2.7, 0.3, size=10, color=MID_GRAY)
    add_txt(sl, val,   x + 0.12, 3.90, 2.7, 0.5, size=13, bold=True, color=DARK_BLUE)

# 底部关键词
kws = ['Carnot Battery', 'Heat Pump + ORC', 'NSGA-II', 'Global Pareto Front',
       'Data Center', 'Conflict Intensity']
for i, kw in enumerate(kws):
    badge(sl, kw, 0.6 + i * 2.1, 5.4,
          color=MID_BLUE if i % 2 == 0 else DARK_BLUE,
          w=1.95, h=0.30)

add_txt(sl, '研究生组会报告', 0.5, 6.85, 12.3, 0.45,
        size=11, color=MID_GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — 研究系统概述
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xF5, 0xF8, 0xFC))
header_bar(sl, '研究系统：Carnot Battery for Data Center',
           'Heat Pump (HP) + Two-tank Storage + Organic Rankine Cycle (ORC)')
slide_num(sl, 2)

# 系统示意文字框
add_rect(sl, 0.3, 1.25, 5.8, 5.8, fill=WHITE, line=RGBColor(0xCC,0xDD,0xEE), line_w=0.8)
add_txt(sl, '系统架构', 0.5, 1.35, 5.4, 0.4, size=14, bold=True, color=DARK_BLUE)

sys_items = [
    ('热泵 (HP)', '将数据中心废热升温至储热罐高温侧', MID_BLUE),
    ('高温储热罐', 'T_st,ht（设计变量，65–140°C）', DARK_BLUE),
    ('低温储热罐', 'T_st,lt = T_cs（与冷源相连）', DARK_BLUE),
    ('有机朗肯循环 (ORC)', '利用储热驱动发电，回收能量', ACCENT_GRN),
]
for i, (name, desc, col) in enumerate(sys_items):
    y = 1.85 + i * 1.15
    add_rect(sl, 0.5, y, 5.3, 0.95, fill=RGBColor(0xE8, 0xF2, 0xFF), line=col, line_w=1.0)
    add_txt(sl, name, 0.65, y + 0.05, 5.0, 0.38, size=13, bold=True, color=col)
    add_txt(sl, desc, 0.65, y + 0.42, 5.0, 0.45, size=11, color=MID_GRAY)

# 四种构型说明
add_rect(sl, 6.3, 1.25, 6.7, 2.65, fill=WHITE, line=RGBColor(0xCC,0xDD,0xEE), line_w=0.8)
add_txt(sl, '四种 CB 构型 (cb_config)', 6.5, 1.35, 6.3, 0.4,
        size=14, bold=True, color=DARK_BLUE)
configs = [
    ('SBVCHP + SBORC', '基础HP + 基础ORC', RGBColor(0x21,0x66,0xAC)),
    ('SRVCHP + SBORC', '回热HP + 基础ORC', RGBColor(0x1A,0x96,0x41)),
    ('SBVCHP + SRORC', '基础HP + 回热ORC', RGBColor(0xD9,0x5F,0x02)),
    ('SRVCHP + SRORC', '回热HP + 回热ORC', RGBColor(0xB2,0x18,0x2B)),
]
for i, (cfg, desc, col) in enumerate(configs):
    y = 1.85 + i * 0.52
    add_rect(sl, 6.5, y, 6.3, 0.44, fill=RGBColor(0xF8,0xF8,0xF8), line=col, line_w=1.2)
    add_txt(sl, cfg,  6.65, y + 0.02, 3.2, 0.38, size=12, bold=True, color=col)
    add_txt(sl, desc, 9.90, y + 0.04, 2.7, 0.35, size=11, color=MID_GRAY)

# 六工况表格
add_rect(sl, 6.3, 4.05, 6.7, 3.0, fill=WHITE, line=RGBColor(0xCC,0xDD,0xEE), line_w=0.8)
add_txt(sl, '六个数据中心工况 (DC-A ~ DC-F)', 6.5, 4.12, 6.3, 0.38,
        size=14, bold=True, color=DARK_BLUE)
headers = ['工况', '场景', 'T_hs', 'T_cs', 'ΔT']
col_x   = [6.45, 7.05, 9.35, 10.35, 11.35]
col_w   = [0.55, 2.25, 0.95, 0.95, 0.85]
for j, (h, x, cw) in enumerate(zip(headers, col_x, col_w)):
    add_rect(sl, x, 4.55, cw, 0.32, fill=DARK_BLUE)
    add_txt(sl, h, x+0.03, 4.57, cw-0.06, 0.28,
            size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
rows = [
    ('DC-A', '空冷·冬', '35°C', '5°C',  '30K'),
    ('DC-B', '空冷·夏', '40°C', '25°C', '15K'),
    ('DC-C', '冷板·冬', '50°C', '5°C',  '45K'),
    ('DC-D', '冷板·夏', '55°C', '25°C', '30K'),
    ('DC-E', '高性液·冬','65°C', '5°C', '60K'),
    ('DC-F', '高性液·夏','75°C', '25°C','50K'),
]
for i, row in enumerate(rows):
    bg = RGBColor(0xE8,0xF2,0xFF) if i % 2 == 0 else WHITE
    for j, (val, x, cw) in enumerate(zip(row, col_x, col_w)):
        add_rect(sl, x, 4.90 + i * 0.33, cw, 0.32, fill=bg)
        col_txt = MID_BLUE if j == 0 else (ACCENT_RED if val.endswith('K') and int(val[:-1])>=45 else BLACK)
        add_txt(sl, val, x+0.03, 4.91+i*0.33, cw-0.06, 0.28,
                size=10, color=col_txt, bold=(j==0),
                align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — 论文框架与当前进度
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xF5, 0xF8, 0xFC))
header_bar(sl, '论文框架与研究进度',
           'Three-part structure: Static MOO → Global Pareto → Off-design Dynamic')
slide_num(sl, 3)

parts = [
    ('Part I', '静态性能映射', '(Mirrors Laterre Ch.3)',
     ['§3  单工况三目标 NSGA-II（DC-A）', '   · H1–H5 五项机制验证 ✓',
      '§4  六工况全局前沿对比', '   · 冲突强度演化图谱 ✓',
      '   · 维度压缩决策（待写作）'],
     ACCENT_GRN, '已完成分析 / 撰写中'),
    ('Part II', '全局前沿与近优设计', '(Mirrors Laterre Ch.4)',
     ['§5  六工况三维全局帕累托面', '   · 三目标非支配排序（待执行）',
      '§6  近优设计空间分析', '   · Must-have vs Real-choice 参数',
      '   · 2–3 个代表设计点提取'],
     MID_BLUE, '规划中'),
    ('Part III', '离设计动态分析', '(ORIGINAL contribution)',
     ['§7  季节边界切换下的性能评估', '   · DC-A↔DC-B, C↔D, E↔F',
      '§8  稳健性分析', '   · 哪种近优设计对季节变化最稳健？',
      '   · 实现 Laterre §3.4.2 的预想'],
     ACCENT_RED, '待执行'),
]

for i, (ptag, ptitle, psub, items, col, status) in enumerate(parts):
    x = 0.25 + i * 4.36
    # 大框
    add_rect(sl, x, 1.15, 4.18, 5.85, fill=WHITE, line=col, line_w=1.5)
    # 顶色条
    add_rect(sl, x, 1.15, 4.18, 0.72, fill=col)
    add_txt(sl, ptag, x + 0.12, 1.18, 0.9, 0.38,
            size=15, bold=True, color=WHITE)
    add_txt(sl, ptitle, x + 1.05, 1.18, 3.0, 0.38,
            size=14, bold=True, color=WHITE)
    add_txt(sl, psub, x + 0.12, 1.55, 3.9, 0.28,
            size=10, color=RGBColor(0xD0,0xE8,0xFF) if col != ACCENT_RED
            else RGBColor(0xFF,0xCC,0xCC))
    # 条目
    for j, item in enumerate(items):
        is_sub = item.startswith('   ')
        fc = MID_GRAY if is_sub else BLACK
        fs = 11 if is_sub else 12
        prefix = '' if is_sub else '▸  '
        done = '✓' in item
        fc = ACCENT_GRN if done else fc
        add_txt(sl, prefix + item.strip(), x + (0.45 if is_sub else 0.2),
                1.98 + j * 0.52, 3.8, 0.48,
                size=fs, color=fc)
    # 状态徽章
    badge_col = ACCENT_GRN if '完成' in status else (MID_BLUE if '规划' in status else MID_GRAY)
    badge(sl, status, x + 0.3, 6.6, color=badge_col, w=3.55, h=0.28)

# 进度箭头
for i in range(2):
    add_txt(sl, '→', 4.38 + i * 4.36, 3.8, 0.5, 0.5,
            size=28, bold=True, color=MID_GRAY, align=PP_ALIGN.CENTER)

# 底部当前位置标注
add_rect(sl, 0.25, 7.1, 12.85, 0.3, fill=RGBColor(0xFF,0xF0,0xC0))
add_txt(sl, '📍 当前位置：Part I §4 六工况全局前沿对比分析（分析完成，正在提炼写作）',
        0.45, 7.12, 12.5, 0.26, size=11, bold=True, color=DARK_BLUE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — 分析方法概述
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xF5, 0xF8, 0xFC))
header_bar(sl, '分析方法：两两全局前沿池化与贡献拆解',
           'Pool all (config × fluid) Pareto solutions → 2D non-dominated sorting → per-pair global front')
slide_num(sl, 4)

# 流程图（横向）
steps = [
    ('42 个流体对\n× 4 种构型\n× 6 个工况',  '数据来源',  MID_BLUE),
    ('Pool 化\n合并所有解',                   '3600–4200解/工况', DARK_BLUE),
    ('三组目标对\n两两非支配排序',             '各得一条全局前沿', MID_BLUE),
    ('构型贡献\n拆解分析',                     '分解图（2×2）', DARK_BLUE),
    ('跨工况\n规律提炼',                       '演化图谱', ACCENT_RED),
]
for i, (title, sub, col) in enumerate(steps):
    x = 0.3 + i * 2.6
    add_rect(sl, x, 1.3, 2.3, 1.5, fill=col)
    add_txt(sl, title, x+0.1, 1.45, 2.1, 0.9,
            size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_txt(sl, sub, x+0.1, 2.3, 2.1, 0.45,
            size=10, color=RGBColor(0xCC,0xDD,0xFF), align=PP_ALIGN.CENTER)
    if i < 4:
        add_txt(sl, '➜', x+2.3, 1.85, 0.3, 0.5,
                size=22, bold=True, color=MID_GRAY, align=PP_ALIGN.CENTER)

# 三组目标对说明
add_rect(sl, 0.3, 3.1, 12.7, 1.55, fill=WHITE, line=RGBColor(0xCC,0xDD,0xEE), line_w=0.8)
add_txt(sl, '三组分析目标对（每个工况各生成 6 张图 = 共 36 张图）',
        0.5, 3.18, 12.3, 0.38, size=13, bold=True, color=DARK_BLUE)
pairs = [
    ('Pair 1', 'η_p2p  vs  e_th', '往返效率 vs 热能量密度', '最强冲突轴（全工况 r = −1.000）', MID_BLUE),
    ('Pair 2', 'η_p2p  vs  η_ex', '往返效率 vs 㶲效率',   'ΔT 依赖型冲突（3.5–39.8 pp 跨度）', DARK_BLUE),
    ('Pair 3', 'η_ex   vs  e_th', '㶲效率 vs 热能量密度', '中等强度冲突（非 Laterre 弱冲突）', ACCENT_RED),
]
for i, (tag, pair, zh, finding, col) in enumerate(pairs):
    x = 0.5 + i * 4.2
    add_rect(sl, x, 3.60, 4.0, 0.95, fill=RGBColor(0xF0,0xF5,0xFF), line=col, line_w=1.0)
    add_txt(sl, f'{tag}:  {pair}', x+0.15, 3.66, 3.7, 0.35, size=12, bold=True, color=col)
    add_txt(sl, zh,      x+0.15, 4.00, 3.7, 0.28, size=10, color=MID_GRAY)
    add_txt(sl, finding, x+0.15, 4.28, 3.7, 0.28, size=10, color=BLACK)

# 图示例说明
add_rect(sl, 0.3, 4.85, 12.7, 2.35, fill=WHITE, line=RGBColor(0xCC,0xDD,0xEE), line_w=0.8)
add_txt(sl, '每对目标输出两张图：',
        0.5, 4.92, 6.0, 0.35, size=13, bold=True, color=DARK_BLUE)
add_txt(sl, '① 总图（Laterre 风格）：个别流体对细线 + 构型 pooled 前沿 + 全局前沿叠加\n'
            '② 分解图（2×2）：每种构型内各流体对的贡献，附全局前沿参考线\n'
            '   ▸ 标注"进入全局前沿的解数 / 总前沿解数"，直读构型支配关系',
        0.5, 5.30, 12.3, 1.75, size=12, color=BLACK)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — DC-A vs DC-E 代表性前沿对比
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xF5, 0xF8, 0xFC))
header_bar(sl, '代表性前沿对比：DC-A（最小 ΔT=30K）vs DC-E（最大 ΔT=60K）',
           'η_p2p–e_th global Pareto front  |  "L-shape plateau" → "uniform ramp"')
slide_num(sl, 5)

add_img(sl, f'{PLOT_DIR}/global_pareto_DCA_laterre_style.png', 0.2, 1.1, 6.3)
add_img(sl, f'{PLOT_DIR}/global_pareto_DCE_p2p_eth_main.png', 6.8, 1.1, 6.3)

# 标注
add_txt(sl, 'DC-A  (ΔT=30K, T_hs=35°C)', 0.35, 1.12, 6.0, 0.38,
        size=12, bold=True, color=MID_BLUE)
add_txt(sl, 'DC-E  (ΔT=60K, T_hs=65°C)', 6.9, 1.12, 6.0, 0.38,
        size=12, bold=True, color=ACCENT_RED)

# 关键对比条
add_rect(sl, 0.2, 6.45, 12.9, 0.85, fill=WHITE, line=MID_BLUE, line_w=0.8)
contrasts = [
    ('η_p2p 上限', '56.9%', '100.0% ⬆'),
    ('e_th 上限',  '34.3 kWh/m³', '45.5 kWh/m³ ⬆'),
    ('前沿形态',   '"宽平台+陡降"L形', '"均匀斜坡"全幅'),
    ('效率极构型', 'SRVCHP_SRORC', 'SRVCHP_SBORC = SRVCHP_SRORC'),
    ('SBVCHP进入η_p2p–η_ex', '否', '是（首次，四构型全入）'),
]
for i, (label, va, vb) in enumerate(contrasts):
    x = 0.35 + i * 2.55
    add_txt(sl, label, x, 6.48, 2.4, 0.28, size=9, color=MID_GRAY, bold=True)
    add_txt(sl, va, x, 6.75, 2.4, 0.28, size=10, color=MID_BLUE)
    add_txt(sl, vb, x, 6.95, 2.4, 0.28, size=10, color=ACCENT_RED, bold=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — 六工况演化图（核心图）
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xF5, 0xF8, 0xFC))
header_bar(sl, '核心图：六工况冲突强度演化图谱',
           'Pareto conflict intensity evolution with ΔT  |  Three objectives, three independent governing mechanisms')
slide_num(sl, 6)

add_img(sl, f'{PLOT_DIR}/six_wp_evolution.png', 0.15, 1.1, 12.95)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — 创新点 1+2（目标上限解耦 + 冲突强度谱）
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xF5, 0xF8, 0xFC))
header_bar(sl, '创新点 C1 + C2：目标上限解耦机制 · 冲突强度谱',
           'Decoupled governing parameters of three objectives  +  ΔT-dependent conflict intensity spectrum')
slide_num(sl, 7)

# C1 框
add_rect(sl, 0.2, 1.15, 6.3, 5.9, fill=WHITE, line=MID_BLUE, line_w=1.5)
add_rect(sl, 0.2, 1.15, 6.3, 0.52, fill=MID_BLUE)
add_txt(sl, 'C1  |  三目标上限的解耦驱动机制', 0.38, 1.18, 5.9, 0.45,
        size=14, bold=True, color=WHITE)
c1_rows = [
    ('η_p2p 上限', '工况 ΔT',
     'ΔT=15K→40%,  30K→57%,  60K→100%\n单调线性，工况决定效率天花板'),
    ('e_th 上限', 'ΔT_sp 允许上限\n（设计约束，非工况）',
     'ΔT_sp≤60K → ~34 kWh/m³\nΔT_sp≤80K → ~45 kWh/m³\n与工况 ΔT 完全无关'),
    ('η_ex 上限', 'ΔT + T_cs\n（双重，非单调）',
     '冬季（T_cs=5°C）随 ΔT 非单调\n夏季（T_cs=25°C）随 ΔT 下降\n不可用单一参数预测'),
]
for i, (obj, driver, evidence) in enumerate(c1_rows):
    y = 1.82 + i * 1.65
    add_rect(sl, 0.35, y, 5.95, 1.5, fill=RGBColor(0xE8,0xF2,0xFF))
    add_txt(sl, obj, 0.5, y + 0.05, 2.2, 0.38, size=13, bold=True, color=DARK_BLUE)
    add_txt(sl, '驱动参数：' + driver, 0.5, y + 0.43, 5.6, 0.55,
            size=11, color=MID_BLUE, bold=True)
    add_txt(sl, evidence, 0.5, y + 0.9, 5.6, 0.55, size=10, color=MID_GRAY)

add_rect(sl, 0.35, 6.7, 5.95, 0.28, fill=YELLOW_HL)
add_txt(sl, '→ 三目标受三类完全不同的参数控制，信息独立，支持三目标优化框架',
        0.5, 6.72, 5.7, 0.24, size=10, bold=True, color=DARK_BLUE)

# C2 框
add_rect(sl, 6.75, 1.15, 6.35, 5.9, fill=WHITE, line=DARK_BLUE, line_w=1.5)
add_rect(sl, 6.75, 1.15, 6.35, 0.52, fill=DARK_BLUE)
add_txt(sl, 'C2  |  ΔT 依赖型冲突强度谱', 6.93, 1.18, 5.9, 0.45,
        size=14, bold=True, color=WHITE)

add_txt(sl, 'η_p2p–η_ex 前沿规模（前沿解数 + 跨度）随 ΔT 的演化：',
        6.93, 1.80, 5.9, 0.4, size=12, color=DARK_BLUE, bold=True)

spectrum = [
    ('DC-B', 15, 9,  3.5, '极弱冲突，几乎可降维',  RGBColor(0xA8,0xC8,0xE8)),
    ('DC-A', 30, 20, 12.0,'弱冲突，降维合理',      RGBColor(0x70,0xA8,0xD8)),
    ('DC-D', 30, 29, 11.5,'弱冲突（高T_hs拓宽）',  RGBColor(0x70,0xA8,0xD8)),
    ('DC-C', 45, 35, 27.1,'中强冲突，降维代价上升', RGBColor(0x38,0x88,0xC8)),
    ('DC-F', 50, 29, 30.0,'中强冲突',              RGBColor(0x20,0x70,0xB8)),
    ('DC-E', 60, 27, 39.8,'强冲突，必须保留三目标', RGBColor(0x10,0x50,0x90)),
]
for i, (wp, dt, n, span, desc, col) in enumerate(spectrum):
    y = 2.28 + i * 0.75
    add_rect(sl, 6.85, y, 6.15, 0.65, fill=col)
    add_txt(sl, f'{wp}  ΔT={dt}K', 7.0, y+0.04, 1.8, 0.28, size=11, bold=True, color=WHITE)
    add_txt(sl, f'n={n}, 跨度={span}pp', 8.85, y+0.04, 1.7, 0.28, size=11, color=WHITE)
    add_txt(sl, desc, 10.6, y+0.04, 2.3, 0.28, size=10, color=YELLOW_HL)
    # 进度条
    bar_w = span / 42.0 * 5.8
    add_rect(sl, 7.0, y+0.38, bar_w, 0.18, fill=WHITE)

add_rect(sl, 6.85, 6.7, 6.15, 0.28, fill=YELLOW_HL)
add_txt(sl, '→ 降维合理性：ΔT ≤ 30K 时可降维；ΔT ≥ 45K 时必须保留三目标',
        7.0, 6.72, 5.9, 0.24, size=10, bold=True, color=DARK_BLUE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — 创新点 3（HP 回热器必要性临界）
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xF5, 0xF8, 0xFC))
header_bar(sl, '创新点 C3：HP 回热器的 ΔT 临界必要性',
           'SBVCHP configs enter η_p2p–η_ex global front only at ΔT ≈ 60 K  |  SBVCHP_SRORC: 17/18 global fronts absent')
slide_num(sl, 8)

# 左：DC-B vs DC-E 的 η_p2p-η_ex 前沿图
add_img(sl, f'{PLOT_DIR}/global_pareto_DCB_p2p_etaex_main.png', 0.2, 1.15, 6.3)
add_img(sl, f'{PLOT_DIR}/global_pareto_DCE_p2p_etaex_main.png', 6.8, 1.15, 6.3)
add_txt(sl, 'DC-B  ΔT=15K — 单构型垄断（SRVCHP_SRORC 100%）',
        0.35, 1.17, 6.2, 0.32, size=11, bold=True, color=MID_BLUE)
add_txt(sl, 'DC-E  ΔT=60K — 四构型全部进入（历史首次）',
        6.95, 1.17, 6.2, 0.32, size=11, bold=True, color=ACCENT_RED)

# 底部规律框
add_rect(sl, 0.2, 6.0, 12.9, 1.3, fill=WHITE, line=ACCENT_RED, line_w=1.5)
add_txt(sl, '跨工况规律：SBVCHP_SRORC 在 18 个全局前沿中 17 次全面缺席',
        0.4, 6.05, 12.5, 0.38, size=14, bold=True, color=ACCENT_RED)

row1 = [('SBVCHP_SRORC（无HP回热 + ORC回热）', '是六工况中最劣构型'),
        ('在 η_p2p–e_th 前沿', '0 次进入（全工况）'),
        ('在 η_ex–e_th 前沿',   '0 次进入（全工况）'),
        ('在 η_p2p–η_ex 前沿', '仅 DC-E 中 3/27（ΔT=60K）')]
for i, (label, val) in enumerate(row1):
    x = 0.4 + i * 3.2
    add_txt(sl, label, x, 6.47, 3.0, 0.28, size=10, color=MID_GRAY)
    add_txt(sl, val, x, 6.75, 3.0, 0.28, size=11, bold=True,
            color=ACCENT_RED if '0' in val else BLACK)

add_rect(sl, 0.4, 7.08, 12.5, 0.18, fill=YELLOW_HL)
add_txt(sl, '工程意义：数据中心场景 ΔT < 60K，HP 回热器具有热力学必要性；ORC 回热器的价值随 ΔT 增大而显著提升',
        0.55, 7.09, 12.2, 0.17, size=10, bold=True, color=DARK_BLUE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — 创新点 4+5（e_th 锁定 + 效率走廊）
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xF5, 0xF8, 0xFC))
header_bar(sl, '创新点 C4 + C5：e_th 设计变量锁定 · 效率走廊跨工况稳定性',
           'Storage density ceiling decoupled from operating condition  |  Efficiency corridor: ΔT_sp → lower bound')
slide_num(sl, 9)

# C4
add_rect(sl, 0.2, 1.15, 6.3, 5.9, fill=WHITE, line=ACCENT_GRN, line_w=1.5)
add_rect(sl, 0.2, 1.15, 6.3, 0.52, fill=ACCENT_GRN)
add_txt(sl, 'C4  |  e_th 天花板的设计变量锁定', 0.38, 1.18, 5.9, 0.45,
        size=14, bold=True, color=WHITE)

add_txt(sl, 'e_th,max 仅由 ΔT_sp 允许上限决定，与工况 ΔT、T_hs、T_cs 完全解耦：',
        0.38, 1.82, 5.9, 0.5, size=12, color=BLACK)
levels = [
    ('ΔT_sp,max = 60 K', '~34 kWh/m³', 'DC-A, DC-B', MID_BLUE),
    ('ΔT_sp,max = 80 K', '~45 kWh/m³', 'DC-C, D, E, F', ACCENT_GRN),
]
for i, (cond, val, wps, col) in enumerate(levels):
    y = 2.45 + i * 1.4
    add_rect(sl, 0.38, y, 5.9, 1.2, fill=RGBColor(0xE8,0xF8,0xEE), line=col, line_w=1.2)
    add_txt(sl, cond, 0.55, y+0.08, 3.0, 0.38, size=14, bold=True, color=col)
    add_txt(sl, '→ ' + val, 0.55, y+0.50, 3.5, 0.38, size=16, bold=True, color=BLACK)
    add_txt(sl, wps, 4.1, y+0.35, 2.0, 0.5, size=11, color=MID_GRAY, align=PP_ALIGN.CENTER)

add_rect(sl, 0.35, 5.42, 5.95, 0.7, fill=YELLOW_HL)
add_txt(sl, '储热密度优化是独立于运行工况的设计约束问题\n→ 可以不依赖工况参数独立优化储热系统规格',
        0.5, 5.45, 5.7, 0.65, size=11, bold=True, color=DARK_BLUE)

add_txt(sl, '注意：DC-A/DC-B 之所以仅能达到 ~34 kWh/m³，\n'
            '是因为储热温差 ΔT_sp 被限制在 60 K 以内，\n'
            '与它们的工况 ΔT 较小（15–30 K）无关。',
        0.38, 6.2, 5.9, 0.9, size=10, color=MID_GRAY, italic=True)

# C5
add_rect(sl, 6.75, 1.15, 6.35, 5.9, fill=WHITE, line=MID_BLUE, line_w=1.5)
add_rect(sl, 6.75, 1.15, 6.35, 0.52, fill=MID_BLUE)
add_txt(sl, 'C5  |  效率走廊的跨工况稳定性', 6.93, 1.18, 5.9, 0.45,
        size=14, bold=True, color=WHITE)

add_txt(sl, '所有工况中，效率极解（高 η_p2p 或高 η_ex）总集中在 ΔT_sp 趋近允许下限处：',
        6.93, 1.82, 5.9, 0.55, size=12, color=BLACK)

corridor_data = [
    ('DC-A / DC-B', 'ΔT_sp,min = 15 K', 'ΔT_sp ≈ 15–18 K', '效率极在 15K 附近'),
    ('DC-C / DC-D', 'ΔT_sp,min = 20 K', 'ΔT_sp ≈ 20–24 K', '效率极在 20K 附近'),
    ('DC-E / DC-F', 'ΔT_sp,min = 25 K', 'ΔT_sp ≈ 25–33 K', '效率极在 25K 附近'),
]
for i, (wps, constraint, actual, desc) in enumerate(corridor_data):
    y = 2.52 + i * 1.25
    add_rect(sl, 6.88, y, 6.08, 1.1, fill=RGBColor(0xE8,0xF2,0xFF))
    add_txt(sl, wps, 7.05, y+0.08, 2.5, 0.3, size=12, bold=True, color=DARK_BLUE)
    add_txt(sl, constraint, 7.05, y+0.40, 2.5, 0.28, size=11, color=MID_GRAY)
    add_txt(sl, actual, 9.65, y+0.12, 2.8, 0.35, size=13, bold=True, color=MID_BLUE)
    add_txt(sl, desc,   9.65, y+0.50, 2.8, 0.28, size=10, color=MID_GRAY)

add_rect(sl, 6.88, 6.28, 6.08, 0.65, fill=YELLOW_HL)
add_txt(sl, '普适设计规律：\n效率优先时令 ΔT_sp → 允许下限（约为下限的 100–130%）\n与工况绝对温度水平无关，可直接用于工程预设计',
        7.0, 6.30, 5.85, 0.63, size=10, bold=True, color=DARK_BLUE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — 与 Laterre 的分歧：维度压缩决策
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xF5, 0xF8, 0xFC))
header_bar(sl, '与 Laterre 论文的关键分歧：维度压缩决策',
           'Laterre: one-cut reduction to 2 objectives  |  Our result: ΔT-dependent, condition-specific')
slide_num(sl, 10)

# 对比表
add_rect(sl, 0.2, 1.15, 12.9, 4.8, fill=WHITE, line=RGBColor(0xCC,0xDD,0xEE), line_w=0.8)
headers = ['分歧点', 'Laterre（原论文）', '本文发现', '依据']
col_xs = [0.25, 2.5, 5.8, 9.8]
col_ws = [2.2, 3.25, 3.95, 3.3]
for j, (h, x, w) in enumerate(zip(headers, col_xs, col_ws)):
    add_rect(sl, x, 1.2, w, 0.38, fill=DARK_BLUE)
    add_txt(sl, h, x+0.08, 1.23, w-0.16, 0.32, size=12, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER)

rows_diff = [
    ('η_II–ρ_el\n(类比 η_ex–e_th)',
     '弱冲突，可降维\n"近共向，信息冗余"',
     '非弱冲突\n六工况均显示实质性冲突（前沿 33–91点）',
     '冲突强度与\nη_p2p–e_th 接近'),
    ('维度压缩策略',
     '一刀切：直接从 3 目标\n降至 2 目标（η_P2P + ρ_el）',
     '工况依赖型：\nΔT ≤ 30K 可降维\nΔT ≥ 45K 必须保留三目标',
     'η_p2p–η_ex 冲突\n强度谱（C2）'),
    ('HP 回热必要性',
     '未系统分析',
     'ΔT < 60K 时 HP 回热\n是热力学必要条件\n（SBVCHP_SRORC 17/18次缺席）',
     '六工况\n构型缺席统计'),
    ('η_p2p 极限',
     '未涉及 100% 场景',
     'DC-E（ΔT=60K）实现\nη_p2p = 100.0%（物理极值）',
     'DC-E 优化结果'),
]
for i, row in enumerate(rows_diff):
    bg = RGBColor(0xF0,0xF5,0xFF) if i % 2 == 0 else WHITE
    for j, (val, x, w) in enumerate(zip(row, col_xs, col_ws)):
        add_rect(sl, x, 1.62 + i * 0.88, w, 0.86, fill=bg)
        col_txt = DARK_BLUE if j == 0 else (ACCENT_RED if j == 2 else BLACK)
        add_txt(sl, val, x+0.1, 1.65+i*0.88, w-0.2, 0.82,
                size=10, color=col_txt, bold=(j == 2))

# 底部框：我们的贡献
add_rect(sl, 0.2, 6.1, 12.9, 1.25, fill=RGBColor(0xE8,0xF5,0xE8), line=ACCENT_GRN, line_w=1.5)
add_txt(sl, '本文对维度压缩问题的原创贡献：',
        0.4, 6.15, 12.5, 0.35, size=13, bold=True, color=ACCENT_GRN)
add_txt(sl, '首次建立"冲突强度谱 → 工况依赖型降维决策"框架，\n'
            '将 Laterre 的静态单工况判断升级为动态多工况判断规则，'
            '直接服务于 Part II 的近优分析目标选择。',
        0.4, 6.52, 12.5, 0.75, size=12, color=DARK_BLUE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — 构型竞争格局演化矩阵
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xF5, 0xF8, 0xFC))
header_bar(sl, '构型竞争格局演化：六工况 × 三目标对',
           'Configuration contribution share in each global Pareto front  |  Color = dominant config')
slide_num(sl, 11)

cfg_colors = {
    'SRVCHP_SBORC': ACCENT_GRN,
    'SRVCHP_SRORC': RGBColor(0xB2,0x18,0x2B),
    'SBVCHP_SBORC': MID_BLUE,
    'SBVCHP_SRORC': RGBColor(0xD9,0x5F,0x02),
    'Mixed':        RGBColor(0x55,0x55,0x88),
}
# 数据：[wp, Pair1_主导, p1_share, Pair2_主导, p2_share, Pair3_主导, p3_share]
matrix = [
    ('DC-B ΔT=15K', 'SRVCHP_SBORC', '66%', 'SRVCHP_SRORC', '100%', 'SRVCHP_SBORC', '79%'),
    ('DC-A ΔT=30K', 'SRVCHP_SBORC', '49%', 'SRVCHP_SRORC', '100%', 'SRVCHP_SBORC', '39%'),
    ('DC-D ΔT=30K', 'SRVCHP_SBORC', '58%', 'SRVCHP_SRORC', '93%', 'SRVCHP_SBORC', '64%'),
    ('DC-C ΔT=45K', 'SRVCHP_SRORC', '57%', 'SRVCHP_SRORC', '83%', 'Mixed',         '50/43%'),
    ('DC-F ΔT=50K', 'SRVCHP_SRORC', '50%', 'SRVCHP_SRORC', '76%', 'SRVCHP_SBORC', '42%'),
    ('DC-E ΔT=60K', 'Mixed',        '47/47%','Mixed',       '全4构型', 'SRVCHP_SBORC','54%'),
]
pair_labels = ['η_p2p–e_th', 'η_p2p–η_ex', 'η_ex–e_th']
col_xs2 = [0.2, 2.8, 6.1, 9.4]
col_ws2 = [2.55, 3.25, 3.25, 3.7]
for j, (h, x, w) in enumerate(zip(['工况'] + pair_labels, col_xs2, col_ws2)):
    add_rect(sl, x, 1.18, w, 0.38, fill=DARK_BLUE)
    add_txt(sl, h, x+0.08, 1.20, w-0.16, 0.34, size=12, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER)

for i, (wp, d1, s1, d2, s2, d3, s3) in enumerate(matrix):
    y = 1.60 + i * 0.84
    bg = RGBColor(0xF0,0xF5,0xFF) if i % 2 == 0 else WHITE
    add_rect(sl, 0.2, y, 2.55, 0.82, fill=bg)
    add_txt(sl, wp, 0.3, y+0.08, 2.35, 0.65,
            size=11, bold=True, color=DARK_BLUE)
    for j, (dom, share) in enumerate([(d1,s1),(d2,s2),(d3,s3)]):
        x = col_xs2[j+1]
        w = col_ws2[j+1]
        cell_col = cfg_colors.get(dom, RGBColor(0x55,0x55,0x88))
        # 饱和度按份额调整
        is_special = '全' in share or '/' in share
        alpha_factor = 0.3 if is_special else float(share.replace('%',''))/100 * 0.5
        add_rect(sl, x+0.05, y+0.05, w-0.1, 0.72, fill=bg)
        add_rect(sl, x+0.05, y+0.05, (w-0.1)*min(float(share.replace('%','').split('/')[0])/100, 1.0)
                 if not is_special else (w-0.1)*0.7, 0.72, fill=cell_col)
        add_txt(sl, dom.replace('_','\n'), x+0.12, y+0.06, w-0.25, 0.45,
                size=9, bold=True, color=WHITE if not is_special else WHITE)
        add_txt(sl, share, x+0.12, y+0.52, w-0.25, 0.24,
                size=10, color=WHITE, bold=True)

# 图例
for i, (cfg, col) in enumerate(cfg_colors.items()):
    badge(sl, cfg.replace('_','+'), 0.3 + i * 2.6, 6.82,
          color=col, w=2.4, h=0.28)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — 当前成果总览
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xF5, 0xF8, 0xFC))
header_bar(sl, '当前成果总览',
           'Deliverables as of 2026-04  |  Analysis complete, writing in progress')
slide_num(sl, 12)

# 左：已完成
add_rect(sl, 0.2, 1.15, 6.3, 5.8, fill=WHITE, line=ACCENT_GRN, line_w=1.5)
add_rect(sl, 0.2, 1.15, 6.3, 0.45, fill=ACCENT_GRN)
add_txt(sl, '✅  已完成', 0.38, 1.18, 5.9, 0.38, size=14, bold=True, color=WHITE)
done_items = [
    ('数据生产', '6 工况 × 42 流体对 × NSGA-II → 23200+ Pareto 解'),
    ('DC-A 精析', '三目标冲突分析，H1–H5 五项机制验证'),
    ('六工况前沿', '18 个两两全局前沿，36 张可视化图'),
    ('五项创新点', 'C1–C5 发掘完毕，数据支撑充分'),
    ('演化图谱', 'six_wp_evolution.png（Figure A 候选）'),
    ('研究报告', '6 份工况报告 + 写作策略笔记（9 个 .md 文件）'),
    ('记忆归档', '关键规律存入 Claude Memory，跨会话可调用'),
]
for i, (tag, desc) in enumerate(done_items):
    add_rect(sl, 0.35, 1.72 + i * 0.68, 5.95, 0.60, fill=RGBColor(0xE8,0xF8,0xEE))
    add_txt(sl, tag,  0.5, 1.75 + i * 0.68, 1.6, 0.28, size=11, bold=True, color=ACCENT_GRN)
    add_txt(sl, desc, 2.2, 1.75 + i * 0.68, 3.9, 0.55, size=10, color=BLACK)

# 右：进行中/待做
add_rect(sl, 6.75, 1.15, 6.35, 5.8, fill=WHITE, line=MID_BLUE, line_w=1.5)
add_rect(sl, 6.75, 1.15, 6.35, 0.45, fill=MID_BLUE)
add_txt(sl, '🔄  进行中 / 待完成', 6.93, 1.18, 5.9, 0.38, size=14, bold=True, color=WHITE)

progress_items = [
    ('撰写中',  MID_BLUE,   'Part I §4 正文（约 1600–2000 字）'),
    ('撰写中',  MID_BLUE,   '维度压缩判断节（与 Laterre 对比辩论）'),
    ('待执行',  MID_GRAY,   '三维全局前沿排序（DC-A + DC-E）'),
    ('待执行',  MID_GRAY,   'Part II §5 三维 Pareto 面分析'),
    ('待执行',  MID_GRAY,   '近优设计空间 + must-have/real-choice 分析'),
    ('待执行',  MID_GRAY,   'Part III 季节对比 / 离设计动态分析'),
    ('规划中',  ACCENT_RED, '构型贡献格局演化热图（Figure B）'),
]
for i, (status, col, desc) in enumerate(progress_items):
    add_rect(sl, 6.88, 1.72 + i * 0.68, 6.08, 0.60, fill=RGBColor(0xF0,0xF5,0xFF))
    badge(sl, status, 6.92, 1.76 + i * 0.68, color=col, w=0.95, h=0.25)
    add_txt(sl, desc, 8.0, 1.75 + i * 0.68, 4.85, 0.55, size=10, color=BLACK)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — 后续研究计划
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xF5, 0xF8, 0xFC))
header_bar(sl, '后续研究与写作计划',
           'Roadmap: Part I completion → Part II execution → Part III dynamic analysis')
slide_num(sl, 13)

phases = [
    ('Phase 1\n（近期）', '完成 Part I 写作',
     ['撰写 §4 正文与维度压缩判断节',
      '制作 Figure A（演化图）、Figure B（热图）',
      '执行 DC-A + DC-E 三维全局前沿排序',
      '完成 Part I 初稿'],
     ACCENT_GRN, '预计 2–3 周'),
    ('Phase 2\n（中期）', 'Part II 全局前沿与近优设计',
     ['六工况三维 Pareto 面计算与可视化',
      '近优设计空间分析（ε-约束法）',
      'Must-have vs Real-choice 参数判断',
      '提取 2–3 个代表设计点'],
     MID_BLUE, '预计 4–6 周'),
    ('Phase 3\n（后期）', 'Part III 季节对比动态分析',
     ['固定设计点，评估季节边界变化',
      'DC-A↔DC-B / DC-C↔DC-D / DC-E↔DC-F',
      '稳健性分析与控制策略建议',
      '完成全文，准备投稿'],
     ACCENT_RED, '预计 6–10 周'),
]

for i, (phase, title, items, col, eta) in enumerate(phases):
    x = 0.25 + i * 4.36
    add_rect(sl, x, 1.15, 4.18, 5.65, fill=WHITE, line=col, line_w=1.5)
    add_rect(sl, x, 1.15, 4.18, 0.72, fill=col)
    add_txt(sl, phase, x+0.12, 1.18, 1.1, 0.65,
            size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_txt(sl, title, x+1.3, 1.30, 2.75, 0.52,
            size=13, bold=True, color=WHITE)
    for j, item in enumerate(items):
        y = 2.02 + j * 1.05
        add_rect(sl, x+0.18, y, 3.8, 0.9, fill=RGBColor(0xF0,0xF5,0xFF))
        add_txt(sl, '▸  ' + item, x+0.3, y+0.08, 3.6, 0.75, size=11, color=BLACK)
    badge(sl, eta, x+0.5, 6.55, color=col, w=3.15, h=0.28)

# 连接箭头
for i in range(2):
    add_txt(sl, '→', 4.38 + i*4.36, 3.5, 0.5, 0.6,
            size=30, bold=True, color=MID_GRAY, align=PP_ALIGN.CENTER)

# 底部时间轴概览
add_rect(sl, 0.2, 7.1, 12.9, 0.32, fill=WHITE, line=MID_GRAY, line_w=0.5)
for i, (label, col) in enumerate([
    ('Part I 初稿', ACCENT_GRN), ('Part II 执行', MID_BLUE),
    ('Part III 执行', ACCENT_RED), ('全文初稿', DARK_BLUE)]):
    x = 0.5 + i * 3.1
    add_rect(sl, x, 7.12, 2.85, 0.26, fill=col)
    add_txt(sl, label, x+0.08, 7.13, 2.7, 0.22,
            size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — 总结
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xF5, 0xF8, 0xFC))
add_rect(sl, 0, 0, 13.33, 1.05, fill=DARK_BLUE)
add_txt(sl, '总结  |  Summary', 0.35, 0.25, 12.5, 0.55,
        size=26, bold=True, color=WHITE)
slide_num(sl, 14)

# 左：五个创新点回顾
add_rect(sl, 0.2, 1.15, 7.5, 5.85, fill=WHITE, line=MID_BLUE, line_w=1.2)
add_txt(sl, '五项核心创新点（已有数据支撑）',
        0.4, 1.22, 7.1, 0.4, size=14, bold=True, color=DARK_BLUE)
innovations = [
    ('C1', '目标上限解耦驱动机制',
     'η_p2p ← ΔT  |  e_th ← ΔT_sp上限  |  η_ex ← ΔT+T_cs（非单调）'),
    ('C2', 'ΔT 依赖型冲突强度谱',
     '前沿跨度从 3.5 pp（DC-B）到 39.8 pp（DC-E），单调演化'),
    ('C3', 'HP 回热器的 ΔT 临界必要性',
     'ΔT < 60K 必须配 HP 回热  |  SBVCHP_SRORC 17/18 次全面缺席'),
    ('C4', 'e_th 天花板的设计变量锁定',
     '与工况完全解耦：ΔT_sp≤60K→34, ≤80K→45 kWh/m³'),
    ('C5', '效率走廊的跨工况稳定性',
     '效率极解恒集中于 ΔT_sp ≈ 允许下限，跨工况普适'),
]
for i, (tag, title, desc) in enumerate(innovations):
    y = 1.72 + i * 1.02
    add_rect(sl, 0.35, y, 7.2, 0.9, fill=RGBColor(0xE8,0xF2,0xFF))
    badge(sl, tag, 0.4, y + 0.28, color=MID_BLUE, w=0.55, h=0.28)
    add_txt(sl, title, 1.05, y + 0.08, 6.3, 0.35, size=12, bold=True, color=DARK_BLUE)
    add_txt(sl, desc,  1.05, y + 0.52, 6.3, 0.35, size=10, color=MID_GRAY)

# 右：关键结论
add_rect(sl, 7.9, 1.15, 5.2, 5.85, fill=WHITE, line=DARK_BLUE, line_w=1.2)
add_txt(sl, '关键结论与后续衔接',
        8.1, 1.22, 4.8, 0.4, size=14, bold=True, color=DARK_BLUE)

conclusions = [
    ('维度压缩', '非一刀切，ΔT 决定\n降维合理性门槛'),
    ('构型筛选', 'SBVCHP_SRORC 可在\n近优分析前排除'),
    ('设计规则', '效率优先→ΔT_sp下限\n密度优先→ΔT_sp上限'),
    ('Part II准备', '三维前沿排序（DC-A\n+ DC-E）下一步执行'),
]
for i, (tag, desc) in enumerate(conclusions):
    y = 1.75 + i * 1.28
    add_rect(sl, 8.05, y, 4.85, 1.12, fill=RGBColor(0xF0,0xF0,0xF8), line=DARK_BLUE, line_w=0.8)
    add_txt(sl, tag,  8.2, y + 0.08, 4.6, 0.32, size=12, bold=True, color=DARK_BLUE)
    add_txt(sl, desc, 8.2, y + 0.45, 4.6, 0.60, size=11, color=BLACK)

# 底部致谢行
add_rect(sl, 0, 7.1, 13.33, 0.4, fill=DARK_BLUE)
add_txt(sl, '感谢聆听  |  欢迎讨论与指导', 0, 7.15, 13.33, 0.32,
        size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ── 保存 ─────────────────────────────────────────────────────────────────────
prs.save(OUT_PATH)
print(f'Saved → {OUT_PATH}')
print(f'Total slides: {len(prs.slides)}')
